"""
Comprehensive tests for document.py — file text extraction.

Tests cover:
- Plain text extraction (txt, md, json, py, etc.)
- PDF extraction (normal and OCR fallback)
- DOCX extraction (paragraphs + tables)
- CSV extraction
- XLSX extraction
- PPTX extraction
- Image OCR
- Error handling and graceful degradation
- Extension normalization
- Preview truncation
"""
import os
import sys
import unittest
from pathlib import Path
from tempfile import NamedTemporaryFile, TemporaryDirectory
from unittest.mock import AsyncMock, MagicMock, patch

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))

# Enable test mode
os.environ["TEST_MODE"] = "1"
from cryptography.fernet import Fernet
_test_key = Fernet.generate_key().decode()
os.environ["MASTER_KEY"] = _test_key

# Mock OCR dependencies BEFORE importing document module
# This ensures OCR_AVAILABLE = True and allows patching of pytesseract/PIL
import types

# Mock pytesseract
pytesseract_mock = types.ModuleType("pytesseract")
pytesseract_mock.image_to_string = MagicMock()
sys.modules["pytesseract"] = pytesseract_mock

# Mock pdf2image
pdf2image_mock = types.ModuleType("pdf2image")
pdf2image_mock.convert_from_path = MagicMock()
sys.modules["pdf2image"] = pdf2image_mock

# Mock PIL and its submodules (pptx imports PIL.ImageFont, etc.)
pil_mock = types.ModuleType("PIL")
pil_mock.__version__ = "10.0.0"  # pypdf checks this
for submodule in ["Image", "ImageFont", "ImageDraw", "ImageFilter", "ImageColor"]:
    sub_mod = types.ModuleType(f"PIL.{submodule}")
    setattr(pil_mock, submodule, sub_mod)
    sys.modules[f"PIL.{submodule}"] = sub_mod
sys.modules["PIL"] = pil_mock

# Now import document - it will see pytesseract and PIL as available
from document import (
    extract_text,
    truncate_preview,
    PLAIN_TEXT_EXTENSIONS,
    IMAGE_EXTENSIONS,
    MAX_OCR_PAGES,
    MAX_OCR_FILE_SIZE_MB,
    OCR_AVAILABLE,
    _extract_pdf,
    _extract_pdf_ocr,
    _extract_docx,
    _extract_csv,
    _extract_xlsx,
    _extract_pptx,
    _extract_image_ocr,
    _extract_plain_text,
)


class PlainTextExtractionTests(unittest.TestCase):
    """Tests for plain text file extraction."""

    def test_extract_txt(self):
        with NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write("Hello, world!")
            f.flush()
            result = extract_text(Path(f.name), "txt")
        self.assertEqual(result, "Hello, world!")

    def test_extract_md(self):
        with NamedTemporaryFile(mode="w", suffix=".md", delete=False) as f:
            f.write("# Title\n\nBody text with **bold**.")
            f.flush()
            result = extract_text(Path(f.name), "md")
        self.assertIn("Title", result)
        self.assertIn("Body text", result)

    def test_extract_json(self):
        with NamedTemporaryFile(mode="w", suffix=".json", delete=False) as f:
            f.write('{"key": "value", "number": 42}')
            f.flush()
            result = extract_text(Path(f.name), "json")
        self.assertIn("key", result)
        self.assertIn("value", result)

    def test_extract_py(self):
        with NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
            f.write("def hello():\n    print('hi')\n")
            f.flush()
            result = extract_text(Path(f.name), "py")
        self.assertIn("def hello", result)

    def test_extract_html(self):
        with NamedTemporaryFile(mode="w", suffix=".html", delete=False) as f:
            f.write("<html><body><p>Hello</p></body></html>")
            f.flush()
            result = extract_text(Path(f.name), "html")
        self.assertIn("Hello", result)

    def test_extract_xml(self):
        with NamedTemporaryFile(mode="w", suffix=".xml", delete=False) as f:
            f.write("<root><child>data</child></root>")
            f.flush()
            result = extract_text(Path(f.name), "xml")
        self.assertIn("data", result)

    def test_extract_rs(self):
        with NamedTemporaryFile(mode="w", suffix=".rs", delete=False) as f:
            f.write("fn main() { println!(\"hi\"); }")
            f.flush()
            result = extract_text(Path(f.name), "rs")
        self.assertIn("fn main", result)

    def test_extract_all_plain_text_extensions(self):
        """Verify all PLAIN_TEXT_EXTENSIONS are supported."""
        for ext in PLAIN_TEXT_EXTENSIONS:
            with NamedTemporaryFile(mode="w", suffix=f".{ext}", delete=False) as f:
                f.write(f"test {ext}")
                f.flush()
                result = extract_text(Path(f.name), ext)
            self.assertIn(ext, result)

    def test_extract_with_utf8_content(self):
        """UTF-8 encoded content is handled correctly."""
        with NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("Hello 世界 🌍")
            f.flush()
            result = extract_text(Path(f.name), "txt")
        self.assertIn("世界", result)
        self.assertIn("🌍", result)


class PDLExtractionTests(unittest.TestCase):
    """Tests for PDF text extraction."""

    def test_extract_pdf_not_found(self):
        """Missing PDF returns error message."""
        result = extract_text(Path("/nonexistent/file.pdf"), "pdf")
        self.assertIn("Could not extract", result)

    def test_extract_pdf_wrong_extension(self):
        """Using .pdf extension on non-PDF returns error."""
        with NamedTemporaryFile(mode="w", suffix=".pdf", delete=False) as f:
            f.write("not a real pdf")
            f.flush()
            result = extract_text(Path(f.name), "pdf")
        # pypdf returns a more specific error about invalid PDF
        self.assertIn("Could not read PDF", result)

    def test_extract_unknown_extension_returns_empty(self):
        """Unknown extension returns empty string."""
        with NamedTemporaryFile(mode="w", suffix=".xyz", delete=False) as f:
            f.write("data")
            f.flush()
            result = extract_text(Path(f.name), "xyz")
        self.assertEqual(result, "")


class DocxExtractionTests(unittest.TestCase):
    """Tests for DOCX extraction."""

    def test_extract_docx_not_found(self):
        result = extract_text(Path("/nonexistent/file.docx"), "docx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_docx(self):
        with NamedTemporaryFile(mode="w", suffix=".docx", delete=False) as f:
            f.write("not a real docx")
            f.flush()
            result = extract_text(Path(f.name), "docx")
        self.assertIn("Could not extract", result)


class CSVExtractionTests(unittest.TestCase):
    """Tests for CSV extraction."""

    def test_extract_csv_not_found(self):
        result = extract_text(Path("/nonexistent/file.csv"), "csv")
        self.assertIn("Could not extract", result)

    def test_extract_empty_csv(self):
        with NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write("")
            f.flush()
            result = extract_text(Path(f.name), "csv")
        self.assertIn("CSV file is empty", result)

    def test_extract_malformed_csv(self):
        # CSV with inconsistent columns - pandas is lenient and parses it
        with NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write('a,b,c\n1,2\n3,')  # inconsistent columns
            f.flush()
            result = extract_text(Path(f.name), "csv")
        # pandas handles this gracefully - verify it returns some content
        self.assertTrue(len(result) > 0)


class XLSXExtractionTests(unittest.TestCase):
    """Tests for XLSX extraction."""

    def test_extract_xlsx_not_found(self):
        result = extract_text(Path("/nonexistent/file.xlsx"), "xlsx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_xlsx(self):
        with NamedTemporaryFile(mode="w", suffix=".xlsx", delete=False) as f:
            f.write("not a real xlsx")
            f.flush()
            result = extract_text(Path(f.name), "xlsx")
        self.assertIn("Could not extract", result)


class PPTXExtractionTests(unittest.TestCase):
    """Tests for PPTX extraction."""

    def test_extract_pptx_not_found(self):
        result = extract_text(Path("/nonexistent/file.pptx"), "pptx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_pptx(self):
        with NamedTemporaryFile(mode="w", suffix=".pptx", delete=False) as f:
            f.write("not a real pptx")
            f.flush()
            result = extract_text(Path(f.name), "pptx")
        self.assertIn("Could not extract", result)


class ImageOCRExtractionTests(unittest.TestCase):
    """Tests for image OCR extraction."""

    def test_extract_image_not_found(self):
        result = extract_text(Path("/nonexistent/file.png"), "png")
        # Could return either file not found error or OCR unavailable message
        self.assertTrue(
            "Could not extract" in result
            or "OCR not available" in result
            or "tesseract" in result.lower()
            or "No text detected" in result
        )

    def test_ocr_handles_missing_tesseract_binary(self):
        """When tesseract binary is missing, extract_text handles it gracefully."""
        import shutil
        # Create a small valid PNG
        png_data = (
            b"\x89PNG\r\n\x1a\n"
            b"\x00\x00\x00\rIHDR"
            b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00"
            b"\x90wS\xde"
            b"\x00\x00\x00\nIDAT"
            b"\x08\xd7c\xf8\xff\xff?\x00\x05\xfe\x02\xfe"
            b"\xa4\xf0\xe3\x18"
            b"\x00\x00\x00\x00IEND"
            b"\xaeB`\x82"
        )
        with NamedTemporaryFile(mode="wb", suffix=".png", delete=False) as f:
            f.write(png_data)
            f.flush()
            result = extract_text(Path(f.name), "png")

        # Should return an error message about missing OCR/tesseract, not crash
        self.assertTrue(
            "OCR not available" in result
            or "Could not extract" in result
            or "tesseract" in result.lower()
            or "No text detected" in result
        )


class ExtensionNormalizationTests(unittest.TestCase):
    """Tests for extension case normalization."""

    def test_uppercase_extension(self):
        with NamedTemporaryFile(mode="w", suffix=".TXT", delete=False) as f:
            f.write("Uppercase extension")
            f.flush()
            result = extract_text(Path(f.name), "TXT")
        self.assertEqual(result, "Uppercase extension")

    def test_mixed_case_extension(self):
        with NamedTemporaryFile(mode="w", suffix=".Md", delete=False) as f:
            f.write("Mixed case extension")
            f.flush()
            result = extract_text(Path(f.name), "Md")
        self.assertEqual(result, "Mixed case extension")

    def test_extension_with_dot(self):
        """Extension with leading dot is normalized."""
        with NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write("With dot")
            f.flush()
            result = extract_text(Path(f.name), ".txt")
        self.assertEqual(result, "With dot")


class TruncatePreviewTests(unittest.TestCase):
    """Tests for truncate_preview function."""

    def test_short_text_unchanged(self):
        self.assertEqual(truncate_preview("Hello"), "Hello")

    def test_long_text_truncated(self):
        text = "a" * 500
        result = truncate_preview(text, length=100)
        self.assertEqual(len(result), 101)  # 100 chars + …
        self.assertTrue(result.endswith("…"))

    def test_exact_length_not_truncated(self):
        text = "a" * 300
        result = truncate_preview(text, length=300)
        self.assertEqual(result, text)

    def test_empty_string_handled(self):
        self.assertEqual(truncate_preview(""), "")

    def test_whitespace_stripped(self):
        result = truncate_preview("  hello  ", length=300)
        self.assertEqual(result, "hello")

    def test_custom_length(self):
        text = "Hello world"
        result = truncate_preview(text, length=5)
        self.assertEqual(result, "Hello…")

    def test_multiline_text(self):
        text = "Line 1\nLine 2\nLine 3"
        result = truncate_preview(text, length=20)
        self.assertTrue(len(result) <= 21)  # 20 + …

    def test_unicode_truncation(self):
        text = "你好世界" * 10  # 40 chars
        result = truncate_preview(text, length=10)
        self.assertTrue(len(result) <= 11)  # 10 + …


class DocxExtractionTests(unittest.TestCase):
    """Tests for DOCX extraction."""

    def test_extract_docx_not_found(self):
        result = extract_text(Path("/nonexistent/file.docx"), "docx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_docx(self):
        with NamedTemporaryFile(mode="w", suffix=".docx", delete=False) as f:
            f.write("not a real docx")
            f.flush()
            result = extract_text(Path(f.name), "docx")
        self.assertIn("Could not extract", result)

    def test_extract_docx_with_tables(self):
        """Test DOCX extraction includes table content."""
        # Create a real DOCX with a table
        from docx import Document as DocxDocument
        from docx.table import Table
        with NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = DocxDocument()
            doc.add_paragraph("Paragraph 1")

            # Add a table
            table = doc.add_table(rows=2, cols=2)
            table.rows[0].cells[0].text = "Header 1"
            table.rows[0].cells[1].text = "Header 2"
            table.rows[1].cells[0].text = "Cell A"
            table.rows[1].cells[1].text = "Cell B"

            doc.add_paragraph("Paragraph 2")
            doc.save(f.name)

            result = extract_text(Path(f.name), "docx")
        self.assertIn("Paragraph 1", result)
        self.assertIn("Paragraph 2", result)
        self.assertIn("Header 1", result)
        self.assertIn("Header 2", result)
        self.assertIn("Cell A", result)
        self.assertIn("Cell B", result)

    def test_extract_docx_multiple_paragraphs(self):
        """Test DOCX with multiple paragraphs."""
        from docx import Document as DocxDocument
        with NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = DocxDocument()
            for i in range(5):
                doc.add_paragraph(f"Paragraph {i}")
            doc.save(f.name)

            result = extract_text(Path(f.name), "docx")
        self.assertIn("Paragraph 0", result)
        self.assertIn("Paragraph 4", result)


class CSVExtractionTests(unittest.TestCase):
    """Tests for CSV extraction."""

    def test_extract_csv_not_found(self):
        result = extract_text(Path("/nonexistent/file.csv"), "csv")
        self.assertIn("Could not extract", result)

    def test_extract_empty_csv(self):
        with NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write("")
            f.flush()
            result = extract_text(Path(f.name), "csv")
        self.assertIn("CSV file is empty", result)

    def test_extract_malformed_csv(self):
        # CSV with inconsistent columns - pandas is lenient and parses it
        with NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write('a,b,c\n1,2\n3,')
            f.flush()
            result = extract_text(Path(f.name), "csv")
        # pandas handles this gracefully - verify it returns some content
        self.assertTrue(len(result) > 0)

    def test_extract_csv_with_data(self):
        with NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write("name,age,city\nJohn,30,NYC\nJane,25,LA\n")
            f.flush()
            result = extract_text(Path(f.name), "csv")
        self.assertIn("name", result)
        self.assertIn("John", result)
        self.assertIn("Jane", result)
        self.assertIn("30", result)

    def test_extract_csv_parser_error(self):
        # Test the ParserError path by creating a file that causes issues
        # This is hard to trigger without specific malformed content
        # Just verify the error handling exists
        pass


class XLSXExtractionTests(unittest.TestCase):
    """Tests for XLSX extraction."""

    def test_extract_xlsx_not_found(self):
        result = extract_text(Path("/nonexistent/file.xlsx"), "xlsx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_xlsx(self):
        with NamedTemporaryFile(mode="w", suffix=".xlsx", delete=False) as f:
            f.write("not a real xlsx")
            f.flush()
            result = extract_text(Path(f.name), "xlsx")
        self.assertIn("Could not extract", result)

    def test_extract_xlsx_multiple_sheets(self):
        """Test XLSX with multiple sheets."""
        import openpyxl
        with NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = openpyxl.Workbook()
            # Sheet 1
            ws1 = wb.active
            ws1.title = "Sheet1"
            ws1.append(["A", "B"])
            ws1.append(["1", "2"])

            # Sheet 2
            ws2 = wb.create_sheet("Sheet2")
            ws2.append(["C", "D"])
            ws2.append(["3", "4"])

            wb.save(f.name)

            result = extract_text(Path(f.name), "xlsx")
        self.assertIn("Sheet: Sheet1", result)
        self.assertIn("Sheet: Sheet2", result)
        self.assertIn("1", result)
        self.assertIn("3", result)

    def test_extract_xlsx_with_none_values(self):
        """Test XLSX with None/empty cells."""
        import openpyxl
        with NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["A", None, "C"])
            ws.append([None, "2", None])
            wb.save(f.name)

            result = extract_text(Path(f.name), "xlsx")
        self.assertIn("", result)  # Empty cells become empty strings


class PPTXExtractionTests(unittest.TestCase):
    """Tests for PPTX extraction."""

    def test_extract_pptx_not_found(self):
        result = extract_text(Path("/nonexistent/file.pptx"), "pptx")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_pptx(self):
        with NamedTemporaryFile(mode="w", suffix=".pptx", delete=False) as f:
            f.write("not a real pptx")
            f.flush()
            result = extract_text(Path(f.name), "pptx")
        self.assertIn("Could not extract", result)

    def test_extract_pptx_multiple_slides_shapes(self):
        """Test PPTX with multiple slides and shapes."""
        from pptx import Presentation
        from pptx.util import Inches
        with NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()

            # Slide 1
            slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
            txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox1.text = "Slide 1 Content"

            # Slide 2
            slide2 = prs.slides.add_slide(prs.slide_layouts[5])
            txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox2.text = "Slide 2 Content"

            prs.save(f.name)

            result = extract_text(Path(f.name), "pptx")
        self.assertIn("Slide 1", result)
        self.assertIn("Slide 1 Content", result)
        self.assertIn("Slide 2", result)
        self.assertIn("Slide 2 Content", result)

    def test_extract_pptx_without_text_shapes(self):
        """Test PPTX with shapes that don't have text frames."""
        from pptx import Presentation
        from pptx.util import Inches
        with NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # Add a shape without text frame (e.g., a rectangle)
            slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))  # MSO_SHAPE.RECTANGLE
            prs.save(f.name)

            result = extract_text(Path(f.name), "pptx")
        self.assertIn("Slide 1", result)


class ImageOCRExtractionTests(unittest.TestCase):
    """Tests for image OCR extraction."""

    def test_extract_image_not_found(self):
        result = extract_text(Path("/nonexistent/file.png"), "png")
        # Could return either file not found error or OCR unavailable message
        self.assertTrue(
            "Could not extract" in result
            or "OCR not available" in result
            or "tesseract" in result.lower()
            or "No text detected" in result
        )

    def test_ocr_handles_missing_tesseract(self):
        """When tesseract binary is missing, returns error message."""
        import shutil
        # Create a small valid PNG
        png_data = (
            b"\x89PNG\r\n\x1a\n"
            b"\x00\x00\x00\rIHDR"
            b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00"
            b"\x90wS\xde"
            b"\x00\x00\x00\nIDAT"
            b"\x08\xd7c\xf8\xff\xff?\x00\x05\xfe\x02\xfe"
            b"\xa4\xf0\xe3\x18"
            b"\x00\x00\x00\x00IEND"
            b"\xaeB`\x82"
        )
        with NamedTemporaryFile(mode="wb", suffix=".png", delete=False) as f:
            f.write(png_data)
            f.flush()
            result = extract_text(Path(f.name), "png")

        # Should return graceful error about OCR/tesseract
        self.assertTrue(
            "OCR not available" in result
            or "Could not extract" in result
            or "tesseract" in result.lower()
            or "No text detected" in result
        )

    def test_image_ocr_size_limit(self):
        """Test image size limit check."""
        # This would require creating a very large image file
        # Just verify the logic exists in the code
        pass


class PDFExtractionTests(unittest.TestCase):
    """Tests for PDF extraction (including OCR fallback)."""

    def test_extract_pdf_not_found(self):
        result = extract_text(Path("/nonexistent/file.pdf"), "pdf")
        self.assertIn("Could not extract", result)

    def test_extract_invalid_pdf(self):
        with NamedTemporaryFile(mode="w", suffix=".pdf", delete=False) as f:
            f.write("not a real pdf")
            f.flush()
            result = extract_text(Path(f.name), "pdf")
        # pypdf raises PdfReadError for invalid PDF
        self.assertIn("Could not read PDF", result)

    @patch("document.OCR_AVAILABLE", True)
    @patch("document._extract_pdf_ocr")
    def test_pdf_ocr_fallback_when_text_sparse(self, mock_ocr):
        """Test OCR fallback when extracted text is too short."""
        mock_ocr.return_value = "OCR extracted text"

        with NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            # Create a minimal valid PDF
            from pypdf import PdfWriter
            writer = PdfWriter()
            writer.add_blank_page(width=100, height=100)
            writer.write(f)
            f.flush()
            result = extract_text(Path(f.name), "pdf")

        # If text is too short (< 100 chars), OCR should be attempted
        mock_ocr.assert_called()

    @patch("document.OCR_AVAILABLE", False)
    def test_pdf_no_ocr_when_unavailable(self):
        """Test PDF extraction doesn't try OCR when unavailable."""
        with NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            from pypdf import PdfWriter
            writer = PdfWriter()
            writer.add_blank_page(width=100, height=100)
            writer.write(f)
            f.flush()
            result = extract_text(Path(f.name), "pdf")
        # Should return the extracted text (empty for blank page)
        self.assertIn("", result)


class PlainTextExtractionInternalTests(unittest.TestCase):
    """Tests for internal plain text extraction function."""

    def test_extract_plain_text_utf8(self):
        with NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("Hello 世界 🌍")
            f.flush()
            result = _extract_plain_text(Path(f.name))
        self.assertIn("世界", result)
        self.assertIn("🌍", result)

    def test_extract_plain_text_encoding_errors_ignored(self):
        # Write with latin-1 but read as utf-8 (will have replacement chars)
        with NamedTemporaryFile(mode="wb", suffix=".txt", delete=False) as f:
            f.write(b"Hello \xe9 world")  # latin-1 e-acute
            f.flush()
            result = _extract_plain_text(Path(f.name))
        self.assertIn("Hello", result)


class InternalFunctionsTests(unittest.TestCase):
    """Tests for internal extraction functions directly."""

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_pdf_ocr_success(self):
        """Test _extract_pdf_ocr with successful conversion."""
        from document import _extract_pdf_ocr

        # Mock pdf2image.convert_from_path inside the function
        with patch("pdf2image.convert_from_path") as mock_convert:
            mock_image = MagicMock()
            mock_convert.return_value = [mock_image]

            # Mock pytesseract
            with patch("document.pytesseract") as mock_pytesseract:
                mock_pytesseract.image_to_string.return_value = "OCR text from page"

                with NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                    from pypdf import PdfWriter
                    writer = PdfWriter()
                    writer.add_blank_page(width=100, height=100)
                    writer.write(f)
                    f.flush()
                    result = _extract_pdf_ocr(Path(f.name))

        self.assertIn("OCR text from page", result)
        self.assertIn("Page 1", result)

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_pdf_ocr_file_too_large(self):
        """Test _extract_pdf_ocr with file exceeding size limit."""
        from document import _extract_pdf_ocr

        with NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"x" * (51 * 1024 * 1024))  # 51 MB > 50 MB limit
            f.flush()
            result = _extract_pdf_ocr(Path(f.name))

        self.assertIn("too large for OCR", result)

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_pdf_ocr_no_pdf2image(self):
        """Test _extract_pdf_ocr when pdf2image not installed."""
        from document import _extract_pdf_ocr

        with patch.dict("sys.modules", {"pdf2image": None}):
            with NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                from pypdf import PdfWriter
                writer = PdfWriter()
                writer.add_blank_page(width=100, height=100)
                writer.write(f)
                f.flush()
                result = _extract_pdf_ocr(Path(f.name))

        self.assertIn("requires pdf2image", result)

    @patch("document.OCR_AVAILABLE", True)
    @patch("document._extract_pdf_ocr")
    def test_extract_pdf_password_protected_returns_password_message(self, mock_ocr):
        """Test password-protected PDF returns password message (doesn't try OCR)."""
        from pypdf.errors import PdfReadError

        with patch("document.PdfReader") as mock_reader:
            mock_reader.side_effect = PdfReadError("password required")

            with NamedTemporaryFile(suffix=".pdf") as f:
                result = _extract_pdf(Path(f.name))

        # Password-protected PDFs return specific message, don't try OCR
        self.assertIn("Password-protected PDF", result)
        mock_ocr.assert_not_called()

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_image_ocr_converts_mode(self):
        """Test _extract_image_ocr converts image mode."""
        from document import _extract_image_ocr

        with patch("document.Image") as mock_image_class:
            mock_image = MagicMock()
            mock_image.mode = "RGBA"  # Not RGB or L
            mock_image.convert.return_value = mock_image
            mock_image_class.open.return_value = mock_image

            with patch("document.pytesseract.image_to_string", return_value="OCR result"):
                with NamedTemporaryFile(suffix=".png") as f:
                    result = _extract_image_ocr(Path(f.name))

        mock_image.convert.assert_called_with("RGB")
        self.assertIn("OCR result", result)

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_image_ocr_no_text_detected(self):
        """Test _extract_image_ocr when no text found."""
        from document import _extract_image_ocr

        with patch("document.Image") as mock_image_class:
            mock_image = MagicMock()
            mock_image.mode = "RGB"
            mock_image_class.open.return_value = mock_image

            with patch("document.pytesseract.image_to_string", return_value="   \n\n  "):
                with NamedTemporaryFile(suffix=".png") as f:
                    result = _extract_image_ocr(Path(f.name))

        self.assertIn("No text detected", result)

    @patch("document.OCR_AVAILABLE", True)
    def test_extract_image_ocr_file_too_large(self):
        """Test _extract_image_ocr with file size limit."""
        from document import _extract_image_ocr

        with NamedTemporaryFile(suffix=".png") as f:
            f.write(b"x" * (51 * 1024 * 1024))  # 51 MB
            f.flush()
            result = _extract_image_ocr(Path(f.name))

        self.assertIn("too large for OCR", result)

    def test_extract_csv_parser_error(self):
        """Test _extract_csv handles ParserError."""
        from document import _extract_csv

        with patch("document.pd.read_csv") as mock_read:
            from pandas.errors import ParserError
            mock_read.side_effect = ParserError("Parse error")

            with NamedTemporaryFile(suffix=".csv") as f:
                result = _extract_csv(Path(f.name))

        self.assertIn("Could not parse CSV", result)

    def test_extract_xlsx_exception(self):
        """Test _extract_xlsx handles general exception."""
        from document import _extract_xlsx

        with patch("document.openpyxl.load_workbook") as mock_load:
            mock_load.side_effect = Exception("XLSX error")

            with NamedTemporaryFile(suffix=".xlsx") as f:
                result = _extract_xlsx(Path(f.name))

        self.assertIn("Could not extract text from this XLSX file", result)

    def test_extract_pptx_exception(self):
        """Test _extract_pptx handles general exception."""
        from document import _extract_pptx

        with patch("document.Presentation") as mock_pres:
            mock_pres.side_effect = Exception("PPTX error")

            with NamedTemporaryFile(suffix=".pptx") as f:
                result = _extract_pptx(Path(f.name))

        self.assertIn("Could not extract text from this PPTX file", result)

    def test_extract_docx_exception(self):
        """Test _extract_docx handles general exception."""
        from document import _extract_docx

        with patch("document.DocxDocument") as mock_docx:
            mock_docx.side_effect = Exception("DOCX error")

            with NamedTemporaryFile(suffix=".docx") as f:
                result = _extract_docx(Path(f.name))

        self.assertIn("Could not extract text from this DOCX file", result)


class ConstantsTests(unittest.TestCase):
    """Tests for module constants."""

    def test_plain_text_extensions_defined(self):
        self.assertIsInstance(PLAIN_TEXT_EXTENSIONS, set)
        self.assertTrue(len(PLAIN_TEXT_EXTENSIONS) > 0)
        # Check common extensions
        for ext in ["txt", "md", "json", "py", "js", "html", "xml", "sql"]:
            self.assertIn(ext, PLAIN_TEXT_EXTENSIONS)

    def test_image_extensions_defined(self):
        self.assertIsInstance(IMAGE_EXTENSIONS, set)
        self.assertTrue(len(IMAGE_EXTENSIONS) > 0)
        for ext in ["png", "jpg", "jpeg", "gif", "webp", "bmp"]:
            self.assertIn(ext, IMAGE_EXTENSIONS)

    def test_constants_positive(self):
        self.assertGreater(MAX_OCR_PAGES, 0)
        self.assertGreater(MAX_OCR_FILE_SIZE_MB, 0)

    def test_ocr_available_type(self):
        self.assertIsInstance(OCR_AVAILABLE, bool)


class ErrorHandlingTests(unittest.TestCase):
    """Tests for error handling in extract_text."""

    def test_permission_error_handled(self):
        """Permission denied returns error message."""
        # We can't easily create a permission-denied file in tests,
        # but we can verify the exception handler exists
        result = extract_text(Path("/root/protected.txt"), "txt")
        self.assertIn("Could not extract", result)

    def test_general_exception_caught(self):
        """Any exception during extraction is caught and returned as message."""
        # The function catches all exceptions and returns error message
        # This is tested implicitly by all the "not found" and "invalid" tests
        pass


if __name__ == "__main__":
    unittest.main()